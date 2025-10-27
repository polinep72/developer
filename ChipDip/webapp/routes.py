# webapp/routes.py
from flask import Blueprint, render_template, request, redirect, url_for, flash, current_app
from flask_login import login_required, current_user
import psycopg2
from psycopg2.extras import RealDictCursor
# Используем модуль database из парсера для некоторых операций
from chipdip_parser import database as db_parser_module
from datetime import datetime, timedelta

# Если формы используются:
# from .forms import AddProductForm

bp = Blueprint('main', __name__)  # Создаем Blueprint


def get_db_connection_webapp():
    """ Получает соединение с БД, используя конфигурацию Flask app """
    db_config = current_app.config['DB_CONFIG']
    try:
        conn = psycopg2.connect(**db_config)
        # current_app.logger.info("DB connection successful for webapp.") # Логирование Flask
        return conn
    except psycopg2.Error as e:
        # current_app.logger.error(f"DB connection error for webapp: {e}")
        flash(f"Ошибка подключения к базе данных: {e}", "error")
        return None


@bp.route('/')
def index():
    # Главная страница с кнопками действий. Списки товаров отображаются на отдельных страницах
    return render_template('index.html')


# Страница со всеми товарами для гостей и авторизованных пользователей
@bp.route('/products')
def all_products():
    conn = get_db_connection_webapp()
    products_data = []
    search_query = request.args.get('q', '').strip()
    if conn:
        try:
            with conn.cursor(cursor_factory=RealDictCursor) as cur:
                if search_query:
                    cur.execute("""
                        SELECT DISTINCT ON (p.url)
                            p.id, p.name, p.internal_sku, p.url, p.current_price, p.last_checked_at,
                            p.user_id,
                            CASE
                                WHEN (SELECT COUNT(*) FROM products p2 WHERE p2.url = p.url AND p2.is_active = TRUE) > 1
                                THEN 'Множественные пользователи'
                                ELSE u.username
                            END as added_by,
                            (SELECT sh.stock_level
                             FROM stock_history sh
                             WHERE sh.product_id = p.id
                             ORDER BY sh.check_timestamp DESC
                             LIMIT 1) as latest_stock
                        FROM products p
                        LEFT JOIN users u ON p.user_id = u.id
                        WHERE p.is_active = TRUE
                          AND LOWER(p.name) LIKE LOWER(%s)
                        ORDER BY p.url, p.id ASC;
                    """, (f"%{search_query}%",))
                else:
                    cur.execute("""
                        SELECT DISTINCT ON (p.url)
                            p.id, p.name, p.internal_sku, p.url, p.current_price, p.last_checked_at,
                            p.user_id, 
                            CASE 
                                WHEN (SELECT COUNT(*) FROM products p2 WHERE p2.url = p.url AND p2.is_active = TRUE) > 1 
                                THEN 'Множественные пользователи'
                                ELSE u.username 
                            END as added_by,
                            (SELECT sh.stock_level 
                             FROM stock_history sh 
                             WHERE sh.product_id = p.id 
                             ORDER BY sh.check_timestamp DESC 
                             LIMIT 1) as latest_stock
                        FROM products p
                        LEFT JOIN users u ON p.user_id = u.id
                        WHERE p.is_active = TRUE
                        ORDER BY p.url, p.id ASC;
                    """)
                products_data = cur.fetchall()
        except psycopg2.Error as e:
            flash(f"Ошибка при загрузке данных: {e}", "error")
        finally:
            conn.close()
    return render_template('products_all.html', products=products_data, q=search_query)


@bp.route('/add_product', methods=['GET', 'POST'])
@login_required
def add_product():
    if request.method == 'POST':
        name = request.form.get('name')  # Это будет internal_sku
        url = request.form.get('url')

        # Простая валидация
        if not name or not url:
            flash('Артикул (Name) и Ссылка (URL) обязательны!', 'error')
            return render_template('add_product.html', name=name, url=url)
        if not url.startswith(('http://', 'https://')):
            flash('URL должен начинаться с http:// или https://', 'error')
            return render_template('add_product.html', name=name, url=url)

        conn = get_db_connection_webapp()
        if conn:
            try:
                with conn.cursor() as cur:
                    # Проверяем, есть ли уже активный товар с таким URL у текущего пользователя
                    cur.execute("SELECT id, is_active FROM products WHERE url = %s AND user_id = %s", (url, current_user.id))
                    existing_product = cur.fetchone()
                    
                    if existing_product:
                        if existing_product[1]:  # is_active = True
                            flash(f"Вы уже добавили товар с URL '{url}'.", "warning")
                        else:
                            # Товар существует у пользователя, но неактивен - активируем его
                            cur.execute(
                                "UPDATE products SET name = %s, internal_sku = %s, is_active = TRUE, updated_at = CURRENT_TIMESTAMP WHERE id = %s",
                                (name, name, existing_product[0])
                            )
                            conn.commit()
                            flash(f"Товар '{name}' успешно восстановлен!", "success")
                            return redirect(url_for('main.my_products'))
                    else:
                        # Используем name из формы как internal_sku, а также как начальное имя товара
                        cur.execute(
                            "INSERT INTO products (name, internal_sku, url, is_active, user_id) VALUES (%s, %s, %s, %s, %s)",
                            (name, name, url, True, current_user.id)
                        )
                        conn.commit()
                        flash(f"Товар '{name}' успешно добавлен!", "success")
                        return redirect(url_for('main.my_products'))  # Перенаправляем в ЛК
            except psycopg2.Error as e:
                conn.rollback()
                flash(f"Ошибка при добавлении товара в БД: {e}", "error")
            finally:
                conn.close()
        else:
            flash("Не удалось подключиться к БД для добавления товара.", "error")

        return render_template('add_product.html', name=name,
                               url=url)  # Возвращаем на форму с данными, если была ошибка

    return render_template('add_product.html')


# Можно добавить маршрут для просмотра истории конкретного товара
@bp.route('/product/<int:product_id>/history')
def product_history(product_id):
    conn = get_db_connection_webapp()
    history_data = []
    product_info = None
    if conn:
        try:
            with conn.cursor(cursor_factory=RealDictCursor) as cur:
                cur.execute("SELECT id, name, url FROM products WHERE id = %s", (product_id,))
                product_info = cur.fetchone()
                if product_info:
                    cur.execute("""
                        SELECT check_timestamp, stock_level, price, raw_text, status, error_message
                        FROM stock_history
                        WHERE product_id = %s
                        ORDER BY check_timestamp DESC
                        LIMIT 100; -- Ограничиваем для производительности
                    """, (product_id,))
                    history_data = cur.fetchall()
                else:
                    flash("Товар не найден.", "error")
                    return redirect(url_for('main.index'))
        except psycopg2.Error as e:
            flash(f"Ошибка при загрузке истории товара: {e}", "error")
        finally:
            conn.close()
    return render_template('product_history.html', product=product_info, history=history_data)

# Личный кабинет - товары пользователя
@bp.route('/my_products')
@login_required
def my_products():
    conn = get_db_connection_webapp()
    products_data = []
    search_query = request.args.get('q', '').strip()
    if conn:
        try:
            with conn.cursor(cursor_factory=RealDictCursor) as cur:
                if search_query:
                    cur.execute("""
                        SELECT 
                            p.id, p.name, p.internal_sku, p.url, p.current_price, p.last_checked_at, p.is_active,
                            (SELECT sh.stock_level 
                             FROM stock_history sh 
                             WHERE sh.product_id = p.id 
                             ORDER BY sh.check_timestamp DESC 
                             LIMIT 1) as latest_stock
                        FROM products p
                        WHERE p.user_id = %s
                          AND LOWER(p.name) LIKE LOWER(%s)
                        ORDER BY p.is_active DESC, p.name ASC;
                    """, (current_user.id, f"%{search_query}%"))
                else:
                    cur.execute("""
                        SELECT 
                            p.id, p.name, p.internal_sku, p.url, p.current_price, p.last_checked_at, p.is_active,
                            (SELECT sh.stock_level 
                             FROM stock_history sh 
                             WHERE sh.product_id = p.id 
                             ORDER BY sh.check_timestamp DESC 
                             LIMIT 1) as latest_stock
                        FROM products p
                        WHERE p.user_id = %s
                        ORDER BY p.is_active DESC, p.name ASC;
                    """, (current_user.id,))
                products_data = cur.fetchall()
        except psycopg2.Error as e:
            flash(f"Ошибка при загрузке данных: {e}", "error")
        finally:
            conn.close()
    return render_template('my_products.html', products=products_data, q=search_query)


# Страница статистики/графиков (заглушка)
@bp.route('/stats')
def stats():
    # Передаем список товаров для селектора на странице статистики
    conn = get_db_connection_webapp()
    products_list = []
    if conn:
        try:
            with conn.cursor(cursor_factory=RealDictCursor) as cur:
                cur.execute("""
                    SELECT id, name, internal_sku
                    FROM products
                    WHERE is_active = TRUE
                    ORDER BY name ASC
                    LIMIT 500
                """)
                products_list = cur.fetchall()
        except psycopg2.Error:
            products_list = []
        finally:
            conn.close()
    return render_template('stats.html', products=products_list)


# --- API: сводная статистика ---
@bp.route('/api/stats/overview')
def api_stats_overview():
    conn = get_db_connection_webapp()
    data = {
        "total_products": 0,
        "active_products": 0,
        "unique_urls": 0,
        "users_count": 0,
        "checks_24h": 0
    }
    if not conn:
        return data
    try:
        with conn.cursor() as cur:
            cur.execute("SELECT COUNT(*) FROM products")
            data["total_products"] = cur.fetchone()[0]

            cur.execute("SELECT COUNT(*) FROM products WHERE is_active = TRUE")
            data["active_products"] = cur.fetchone()[0]

            cur.execute("SELECT COUNT(DISTINCT url) FROM products WHERE is_active = TRUE")
            data["unique_urls"] = cur.fetchone()[0]

            cur.execute("SELECT COUNT(*) FROM users WHERE is_active = TRUE")
            data["users_count"] = cur.fetchone()[0]

            cur.execute("""
                SELECT COUNT(*)
                FROM stock_history
                WHERE check_timestamp >= NOW() - INTERVAL '24 hours'
            """)
            data["checks_24h"] = cur.fetchone()[0]
    finally:
        conn.close()
    return data


# --- API: временной ряд по товару ---
@bp.route('/api/stats/product/<int:product_id>/series')
def api_stats_product_series(product_id: int):
    conn = get_db_connection_webapp()
    if not conn:
        return {"points": []}
    points = []
    try:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute(
                """
                SELECT to_char(check_timestamp, 'YYYY-MM-DD HH24:MI') AS ts,
                       stock_level,
                       price
                FROM stock_history
                WHERE product_id = %s
                ORDER BY check_timestamp ASC
                LIMIT 500
                """,
                (product_id,)
            )
            for row in cur.fetchall():
                points.append({
                    "ts": row["ts"],
                    "stock": row["stock_level"],
                    "price": float(row["price"]) if row["price"] is not None else None
                })
    finally:
        conn.close()
    return {"points": points}


# --- API: продажи по убыванию остатков за интервал ---
@bp.route('/api/stats/sales')
def api_stats_sales():
    # Параметры интервала
    date_from_str = request.args.get('from')
    date_to_str = request.args.get('to')
    now_dt = datetime.now()
    if not date_to_str:
        date_to = now_dt
    else:
        # Разрешаем формат YYYY-MM-DD
        try:
            date_to = datetime.strptime(date_to_str, '%Y-%m-%d') + timedelta(days=1)
        except ValueError:
            return {"error": "Неверный формат 'to' (используйте YYYY-MM-DD)"}, 400
    if not date_from_str:
        date_from = now_dt - timedelta(days=7)
    else:
        try:
            date_from = datetime.strptime(date_from_str, '%Y-%m-%d')
        except ValueError:
            return {"error": "Неверный формат 'from' (используйте YYYY-MM-DD)"}, 400

    # Гарантия: from < to
    if date_from >= date_to:
        return {"error": "Параметр 'from' должен быть раньше 'to'"}, 400

    conn = get_db_connection_webapp()
    if not conn:
        return {"items": []}
    items = []
    try:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute(
                """
                WITH sh AS (
                  SELECT
                    p.id AS product_id,
                    p.name,
                    p.internal_sku,
                    sh.check_timestamp,
                    sh.stock_level,
                    LAG(sh.stock_level) OVER (PARTITION BY sh.product_id ORDER BY sh.check_timestamp) AS prev_stock
                  FROM stock_history sh
                  JOIN products p ON p.id = sh.product_id
                  WHERE COALESCE(p.is_active, TRUE) = TRUE
                    AND sh.stock_level IS NOT NULL
                ),
                sales_calc AS (
                  SELECT
                    product_id,
                    name,
                    internal_sku,
                    check_timestamp,
                    stock_level,
                    prev_stock,
                    -- Рассчитываем продажи только если текущий остаток меньше предыдущего
                    -- (исключаем пополнения склада)
                    CASE 
                      WHEN prev_stock IS NOT NULL AND stock_level < prev_stock 
                      THEN prev_stock - stock_level
                      ELSE 0
                    END AS sold_qty
                  FROM sh
                  WHERE check_timestamp >= %s AND check_timestamp < %s
                )
                SELECT
                  product_id,
                  name,
                  internal_sku,
                  SUM(sold_qty)::bigint AS sold_qty
                FROM sales_calc
                GROUP BY product_id, name, internal_sku
                HAVING SUM(sold_qty) > 0
                ORDER BY sold_qty DESC, name ASC
                LIMIT 100
                """,
                (date_from, date_to)
            )
            rows = cur.fetchall()
            for r in rows:
                items.append({
                    "product_id": r["product_id"],
                    "name": r["name"],
                    "internal_sku": r["internal_sku"],
                    "sold_qty": int(r["sold_qty"]) if r["sold_qty"] is not None else 0
                })
    finally:
        conn.close()
    return {"items": items, "from": date_from.strftime('%Y-%m-%d'), "to": (date_to - timedelta(days=1)).strftime('%Y-%m-%d')}


# --- API: годовые продажи по месяцам для всех товаров ---
@bp.route('/api/stats/yearly/monthly')
def api_stats_yearly_monthly():
    """API для получения продаж по месяцам за последние 12 месяцев для всех товаров"""
    conn = get_db_connection_webapp()
    if not conn:
        return {"months": [], "products": []}
    months = []
    products = []
    try:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            # Получаем все месяцы
            cur.execute("""
                SELECT 
                  generate_series(
                    DATE_TRUNC('month', CURRENT_DATE - INTERVAL '11 months'),
                    DATE_TRUNC('month', CURRENT_DATE),
                    INTERVAL '1 month'
                  )::date AS month_start
                ORDER BY month_start
            """)
            
            for row in cur.fetchall():
                months.append(row['month_start'].strftime('%b %Y'))
            
            # Получаем продажи по товарам и месяцам (исправленная логика)
            cur.execute("""
                WITH sh AS (
                  SELECT
                    p.id AS product_id,
                    p.name as product_name,
                    sh.check_timestamp,
                    sh.stock_level,
                    LAG(sh.stock_level) OVER (PARTITION BY sh.product_id ORDER BY sh.check_timestamp) AS prev_stock
                  FROM stock_history sh
                  JOIN products p ON p.id = sh.product_id
                  WHERE COALESCE(p.is_active, TRUE) = TRUE
                    AND sh.stock_level IS NOT NULL
                    AND sh.check_timestamp >= DATE_TRUNC('month', CURRENT_DATE - INTERVAL '11 months')
                ),
                sales_calc AS (
                  SELECT
                    product_id,
                    product_name,
                    DATE_TRUNC('month', check_timestamp) AS month_start,
                    SUM(
                      CASE 
                        WHEN prev_stock IS NOT NULL AND stock_level < prev_stock 
                        THEN prev_stock - stock_level
                        ELSE 0
                      END
                    ) AS sold_qty
                  FROM sh
                  GROUP BY product_id, product_name, DATE_TRUNC('month', check_timestamp)
                  HAVING SUM(
                    CASE 
                      WHEN prev_stock IS NOT NULL AND stock_level < prev_stock 
                      THEN prev_stock - stock_level
                      ELSE 0
                    END
                  ) > 0
                )
                SELECT 
                  product_id,
                  product_name,
                  month_start,
                  sold_qty
                FROM sales_calc
                ORDER BY product_id, month_start
            """)
            
            # Группируем данные по товарам
            product_data = {}
            for row in cur.fetchall():
                product_id = row['product_id']
                if product_id not in product_data:
                    product_data[product_id] = {
                        'id': product_id,
                        'name': row['product_name'],
                        'sales': [0] * len(months)
                    }
                
                month_name = row['month_start'].strftime('%b %Y')
                month_index = months.index(month_name)
                product_data[product_id]['sales'][month_index] = int(row['sold_qty'])
            
            products = list(product_data.values())
            
    except Exception as e:
        current_app.logger.error(f"Ошибка получения данных для годового графика: {e}", exc_info=True)
    finally:
        conn.close()
    return {"months": months, "products": products}

# --- API: месячные продажи по товару ---
@bp.route('/api/stats/product/<int:product_id>/monthly')
def api_stats_product_monthly(product_id: int):
    conn = get_db_connection_webapp()
    if not conn:
        return {"points": []}
    points = []
    try:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute(
                """
                WITH months_series AS (
                  -- Генерируем последние 12 месяцев
                  SELECT 
                    DATE_TRUNC('month', CURRENT_DATE - INTERVAL '11 months' + (generate_series(0, 11) * INTERVAL '1 month')) AS month_start
                ),
                sh AS (
                  SELECT
                    sh.check_timestamp,
                    sh.stock_level,
                    LAG(sh.stock_level) OVER (ORDER BY sh.check_timestamp) AS prev_stock
                  FROM stock_history sh
                  WHERE sh.product_id = %s
                    AND sh.stock_level IS NOT NULL
                    AND sh.check_timestamp >= DATE_TRUNC('month', CURRENT_DATE - INTERVAL '11 months')
                ),
                sales_calc AS (
                  SELECT
                    DATE_TRUNC('month', check_timestamp) AS month_start,
                    SUM(
                      CASE 
                        WHEN prev_stock IS NOT NULL AND stock_level < prev_stock 
                        THEN prev_stock - stock_level
                        ELSE 0
                      END
                    ) AS sold_qty
                  FROM sh
                  GROUP BY DATE_TRUNC('month', check_timestamp)
                  HAVING SUM(
                    CASE 
                      WHEN prev_stock IS NOT NULL AND stock_level < prev_stock 
                      THEN prev_stock - stock_level
                      ELSE 0
                    END
                  ) > 0  -- Исключаем месяцы без продаж
                )
                SELECT 
                  ms.month_start,
                  COALESCE(sc.sold_qty, 0) AS sold_qty
                FROM months_series ms
                LEFT JOIN sales_calc sc ON ms.month_start = sc.month_start
                ORDER BY ms.month_start ASC
                """,
                (product_id,)
            )
            for row in cur.fetchall():
                points.append({
                    "month": row["month_start"].strftime('%Y-%m'),
                    "sold_qty": int(row["sold_qty"]) if row["sold_qty"] is not None else 0
                })
    finally:
        conn.close()
    return {"points": points}

# Удаление товара пользователем
@bp.route('/delete_product/<int:product_id>', methods=['POST'])
@login_required
def delete_product(product_id):
    conn = get_db_connection_webapp()
    if conn:
        try:
            with conn.cursor() as cur:
                # Проверяем, принадлежит ли товар текущему пользователю
                cur.execute("SELECT id FROM products WHERE id = %s AND user_id = %s", (product_id, current_user.id))
                if not cur.fetchone():
                    flash("Товар не найден или у вас нет прав для его удаления.", "error")
                    return redirect(url_for('main.my_products'))
                
                # Удаляем товар (мягкое удаление - устанавливаем is_active = FALSE)
                cur.execute("UPDATE products SET is_active = FALSE WHERE id = %s", (product_id,))
                conn.commit()
                flash("Товар успешно удален.", "success")
        except psycopg2.Error as e:
            conn.rollback()
            flash(f"Ошибка при удалении товара: {e}", "error")
        finally:
            conn.close()
    else:
        flash("Не удалось подключиться к базе данных.", "error")
    
    return redirect(url_for('main.my_products'))

# Восстановление товара
@bp.route('/restore_product/<int:product_id>', methods=['POST'])
@login_required
def restore_product(product_id):
    conn = get_db_connection_webapp()
    if conn:
        try:
            with conn.cursor() as cur:
                # Проверяем, принадлежит ли товар текущему пользователю и неактивен ли он
                cur.execute("SELECT id FROM products WHERE id = %s AND user_id = %s AND is_active = FALSE", (product_id, current_user.id))
                if not cur.fetchone():
                    flash("Товар не найден или у вас нет прав для его восстановления.", "error")
                    return redirect(url_for('main.my_products'))
                
                # Восстанавливаем товар
                cur.execute("UPDATE products SET is_active = TRUE, updated_at = CURRENT_TIMESTAMP WHERE id = %s", (product_id,))
                conn.commit()
                flash("Товар успешно восстановлен.", "success")
        except psycopg2.Error as e:
            conn.rollback()
            flash(f"Ошибка при восстановлении товара: {e}", "error")
        finally:
            conn.close()
    else:
        flash("Не удалось подключиться к базе данных.", "error")
    
    return redirect(url_for('main.my_products'))


# --- API: обновление названия товара ---
@bp.route('/api/product/<int:product_id>/name', methods=['PUT'])
@login_required
def api_update_product_name(product_id: int):
    """API для обновления названия товара (только для владельца товара)"""
    conn = get_db_connection_webapp()
    if not conn:
        return {"success": False, "error": "Ошибка подключения к БД"}, 500
    
    try:
        data = request.get_json()
        new_name = data.get('name', '').strip()
        
        if not new_name:
            return {"success": False, "error": "Название не может быть пустым"}, 400
        
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            # Проверяем, что товар принадлежит текущему пользователю (включая удаленные)
            cur.execute("""
                SELECT id, name, user_id, is_active FROM products 
                WHERE id = %s
            """, (product_id,))
            
            product = cur.fetchone()
            current_app.logger.info(f"Debug: product_id={product_id}, current_user.id={current_user.id}, product={product}")
            
            if not product:
                return {"success": False, "error": "Товар не найден"}, 404
            
            if product['user_id'] != current_user.id:
                return {"success": False, "error": "У вас нет прав на редактирование этого товара"}, 403
            
            # Обновляем название
            cur.execute("""
                UPDATE products 
                SET name = %s, updated_at = NOW() 
                WHERE id = %s
            """, (new_name, product_id))
            
            conn.commit()
            current_app.logger.info(f"Пользователь {current_user.id} обновил название товара {product_id}: '{product['name']}' -> '{new_name}'")
            
            return {"success": True, "new_name": new_name}
            
    except Exception as e:
        conn.rollback()
        current_app.logger.error(f"Ошибка обновления названия товара {product_id}: {e}", exc_info=True)
        return {"success": False, "error": "Ошибка при обновлении названия"}, 500
    finally:
        conn.close()


# --- API: экспорт данных в Excel ---
@bp.route('/api/export/excel')
@login_required
def api_export_excel():
    """API для экспорта данных графиков в Excel"""
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, LineChart, Reference
    from openpyxl.chart.axis import DateAxis
    from io import BytesIO
    import tempfile
    import os
    
    try:
        # Создаем Excel файл
        wb = Workbook()
        ws = wb.active
        ws.title = "Статистика ChipDip"
        
        # Получаем данные для экспорта
        conn = get_db_connection_webapp()
        if not conn:
            return {"error": "Ошибка подключения к БД"}, 500
        
        # Добавляем заголовки
        ws['A1'] = "Статистика продаж ChipDip"
        ws['A2'] = "Дата экспорта: " + datetime.now().strftime('%d.%m.%Y %H:%M')
        
        # Получаем данные продаж за последние 12 месяцев
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                WITH months_series AS (
                  SELECT 
                    DATE_TRUNC('month', CURRENT_DATE - INTERVAL '11 months' + (generate_series(0, 11) * INTERVAL '1 month')) AS month_start
                ),
                sh AS (
                  SELECT
                    p.id AS product_id,
                    p.name as product_name,
                    sh.check_timestamp,
                    sh.stock_level,
                    LAG(sh.stock_level) OVER (PARTITION BY sh.product_id ORDER BY sh.check_timestamp) AS prev_stock
                  FROM stock_history sh
                  JOIN products p ON p.id = sh.product_id
                  WHERE COALESCE(p.is_active, TRUE) = TRUE
                    AND sh.stock_level IS NOT NULL
                    AND sh.check_timestamp >= DATE_TRUNC('month', CURRENT_DATE - INTERVAL '11 months')
                ),
                sales_calc AS (
                  SELECT
                    product_id,
                    product_name,
                    DATE_TRUNC('month', check_timestamp) AS month_start,
                    SUM(
                      CASE 
                        WHEN prev_stock IS NOT NULL AND stock_level < prev_stock 
                        THEN prev_stock - stock_level
                        ELSE 0
                      END
                    ) AS sold_qty
                  FROM sh
                  GROUP BY product_id, product_name, DATE_TRUNC('month', check_timestamp)
                  HAVING SUM(
                    CASE 
                      WHEN prev_stock IS NOT NULL AND stock_level < prev_stock 
                      THEN prev_stock - stock_level
                      ELSE 0
                    END
                  ) > 0
                )
                SELECT 
                  product_name,
                  month_start,
                  sold_qty
                FROM sales_calc
                ORDER BY product_name, month_start
            """)
            
            # Записываем данные в Excel
            row = 4
            ws['A4'] = "Товар"
            ws['B4'] = "Месяц"
            ws['C4'] = "Продано (шт)"
            
            for row_data in cur.fetchall():
                row += 1
                ws[f'A{row}'] = row_data['product_name']
                ws[f'B{row}'] = row_data['month_start'].strftime('%Y-%m')
                ws[f'C{row}'] = int(row_data['sold_qty'])
        
        conn.close()
        
        # Создаем временный файл
        with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp_file:
            wb.save(tmp_file.name)
            
            # Читаем файл и отправляем пользователю
            with open(tmp_file.name, 'rb') as f:
                file_data = f.read()
            
            # Удаляем временный файл
            os.unlink(tmp_file.name)
            
            from flask import Response
            return Response(
                file_data,
                mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                headers={
                    'Content-Disposition': f'attachment; filename=chipdip_stats_{datetime.now().strftime("%Y%m%d_%H%M")}.xlsx'
                }
            )
            
    except Exception as e:
        current_app.logger.error(f"Ошибка экспорта в Excel: {e}", exc_info=True)
        return {"error": "Ошибка при создании Excel файла"}, 500


# --- API: экспорт данных в PowerPoint ---
@bp.route('/api/export/powerpoint')
@login_required
def api_export_powerpoint():
    """API для экспорта данных графиков в PowerPoint"""
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from io import BytesIO
    import tempfile
    import os
    
    try:
        # Создаем PowerPoint презентацию
        prs = Presentation()
        
        # Добавляем титульный слайд
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Статистика ChipDip"
        subtitle.text = f"Отчет за {datetime.now().strftime('%d.%m.%Y %H:%M')}"
        
        # Получаем данные для экспорта
        conn = get_db_connection_webapp()
        if not conn:
            return {"error": "Ошибка подключения к БД"}, 500
        
        # Добавляем слайд с данными
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        title.text = "Статистика продаж"
        
        # Получаем данные продаж
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                WITH months_series AS (
                  SELECT 
                    DATE_TRUNC('month', CURRENT_DATE - INTERVAL '11 months' + (generate_series(0, 11) * INTERVAL '1 month')) AS month_start
                ),
                sh AS (
                  SELECT
                    p.id AS product_id,
                    p.name as product_name,
                    sh.check_timestamp,
                    sh.stock_level,
                    LAG(sh.stock_level) OVER (PARTITION BY sh.product_id ORDER BY sh.check_timestamp) AS prev_stock
                  FROM stock_history sh
                  JOIN products p ON p.id = sh.product_id
                  WHERE COALESCE(p.is_active, TRUE) = TRUE
                    AND sh.stock_level IS NOT NULL
                    AND sh.check_timestamp >= DATE_TRUNC('month', CURRENT_DATE - INTERVAL '11 months')
                ),
                sales_calc AS (
                  SELECT
                    product_id,
                    product_name,
                    DATE_TRUNC('month', check_timestamp) AS month_start,
                    SUM(
                      CASE 
                        WHEN prev_stock IS NOT NULL AND stock_level < prev_stock 
                        THEN prev_stock - stock_level
                        ELSE 0
                      END
                    ) AS sold_qty
                  FROM sh
                  GROUP BY product_id, product_name, DATE_TRUNC('month', check_timestamp)
                  HAVING SUM(
                    CASE 
                      WHEN prev_stock IS NOT NULL AND stock_level < prev_stock 
                      THEN prev_stock - stock_level
                      ELSE 0
                    END
                  ) > 0
                )
                SELECT 
                  product_name,
                  month_start,
                  sold_qty
                FROM sales_calc
                ORDER BY product_name, month_start
            """)
            
            # Создаем таблицу с данными
            rows = cur.fetchall()
            if rows:
                # Создаем таблицу
                table_data = [["Товар", "Месяц", "Продано (шт)"]]
                for row_data in rows:
                    table_data.append([
                        row_data['product_name'],
                        row_data['month_start'].strftime('%Y-%m'),
                        str(int(row_data['sold_qty']))
                    ])
                
                # Добавляем таблицу на слайд
                rows_count = len(table_data)
                cols_count = len(table_data[0])
                
                table = slide.shapes.add_table(rows_count, cols_count, 100, 200, 600, 300).table
                
                for i, row_data in enumerate(table_data):
                    for j, cell_data in enumerate(row_data):
                        table.cell(i, j).text = cell_data
        
        conn.close()
        
        # Создаем временный файл
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
            prs.save(tmp_file.name)
            
            # Читаем файл и отправляем пользователю
            with open(tmp_file.name, 'rb') as f:
                file_data = f.read()
            
            # Удаляем временный файл
            os.unlink(tmp_file.name)
            
            from flask import Response
            return Response(
                file_data,
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                headers={
                    'Content-Disposition': f'attachment; filename=chipdip_stats_{datetime.now().strftime("%Y%m%d_%H%M")}.pptx'
                }
            )
            
    except Exception as e:
        current_app.logger.error(f"Ошибка экспорта в PowerPoint: {e}", exc_info=True)
        return {"error": "Ошибка при создании PowerPoint файла"}, 500


# --- API: экспорт графика продаж по убыванию в Excel ---
@bp.route('/api/export/chart/sales/excel')
@login_required
def api_export_chart_sales_excel():
    """API для экспорта графика продаж по убыванию в Excel"""
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, Reference
    from io import BytesIO
    import tempfile
    import os

    try:
        # Создаем Excel файл
        wb = Workbook()
        ws = wb.active
        ws.title = "Продажи по убыванию"

        # Получаем данные для экспорта
        conn = get_db_connection_webapp()
        if not conn:
            return {"error": "Ошибка подключения к БД"}, 500

        # Добавляем заголовки
        ws['A1'] = "График продаж по убыванию"
        ws['A2'] = "Дата экспорта: " + datetime.now().strftime('%d.%m.%Y %H:%M')

        # Получаем данные продаж за последние 30 дней
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                WITH sh AS (
                  SELECT
                    p.id AS product_id,
                    p.name as product_name,
                    sh.check_timestamp,
                    sh.stock_level,
                    LAG(sh.stock_level) OVER (PARTITION BY sh.product_id ORDER BY sh.check_timestamp) AS prev_stock
                  FROM stock_history sh
                  JOIN products p ON p.id = sh.product_id
                  WHERE COALESCE(p.is_active, TRUE) = TRUE
                    AND sh.stock_level IS NOT NULL
                    AND sh.check_timestamp >= CURRENT_DATE - INTERVAL '30 days'
                ),
                sales_calc AS (
                  SELECT
                    product_id,
                    product_name,
                    SUM(
                      CASE
                        WHEN prev_stock IS NOT NULL AND stock_level < prev_stock
                        THEN prev_stock - stock_level
                        ELSE 0
                      END
                    ) AS sold_qty
                  FROM sh
                  GROUP BY product_id, product_name
                  HAVING SUM(
                    CASE
                      WHEN prev_stock IS NOT NULL AND stock_level < prev_stock
                      THEN prev_stock - stock_level
                      ELSE 0
                    END
                  ) > 0
                )
                SELECT
                  product_name,
                  sold_qty
                FROM sales_calc
                ORDER BY sold_qty DESC
                LIMIT 20
            """)

            # Записываем данные в Excel
            row = 4
            ws['A4'] = "Товар"
            ws['B4'] = "Продано (шт)"

            for row_data in cur.fetchall():
                row += 1
                ws[f'A{row}'] = row_data['product_name']
                ws[f'B{row}'] = int(row_data['sold_qty'])

        conn.close()

        # Создаем график
        chart = BarChart()
        chart.title = "Топ-20 товаров по продажам"
        chart.style = 10
        chart.y_axis.title = 'Количество проданных товаров'
        chart.x_axis.title = 'Товар'

        # Добавляем данные для графика
        data = Reference(ws, min_col=2, min_row=4, max_row=row)
        cats = Reference(ws, min_col=1, min_row=5, max_row=row)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)

        # Добавляем график на лист
        ws.add_chart(chart, "D2")

        # Создаем временный файл
        with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp_file:
            wb.save(tmp_file.name)

            # Читаем файл и отправляем пользователю
            with open(tmp_file.name, 'rb') as f:
                file_data = f.read()

            # Удаляем временный файл
            os.unlink(tmp_file.name)

            from flask import Response
            return Response(
                file_data,
                mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                headers={
                    'Content-Disposition': f'attachment; filename=chipdip_sales_{datetime.now().strftime("%Y%m%d_%H%M")}.xlsx'
                }
            )

    except Exception as e:
        current_app.logger.error(f"Ошибка экспорта графика продаж в Excel: {e}", exc_info=True)
        return {"error": "Ошибка при создании Excel файла с графиком продаж"}, 500


# --- API: статистика выручки по месяцам (общий) ---
@bp.route('/api/stats/revenue/monthly')
@login_required
def api_stats_revenue_monthly():
    """API для получения данных выручки по месяцам для всех товаров"""
    conn = get_db_connection_webapp()
    if not conn:
        return {"error": "Ошибка подключения к БД"}, 500

    try:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                WITH sh AS (
                  SELECT
                    p.id AS product_id,
                    p.name as product_name,
                    sh.check_timestamp,
                    sh.stock_level,
                    sh.price,
                    LAG(sh.stock_level) OVER (PARTITION BY sh.product_id ORDER BY sh.check_timestamp) AS prev_stock
                  FROM stock_history sh
                  JOIN products p ON p.id = sh.product_id
                  WHERE COALESCE(p.is_active, TRUE) = TRUE
                    AND sh.stock_level IS NOT NULL
                    AND sh.price IS NOT NULL
                    AND sh.check_timestamp >= DATE_TRUNC('month', CURRENT_DATE - INTERVAL '11 months')
                ),
                sales_calc AS (
                  SELECT
                    product_id,
                    product_name,
                    DATE_TRUNC('month', check_timestamp) AS month_start,
                    SUM(
                      CASE
                        WHEN prev_stock IS NOT NULL AND stock_level < prev_stock
                        THEN (prev_stock - stock_level) * price
                        ELSE 0
                      END
                    ) AS revenue
                  FROM sh
                  GROUP BY product_id, product_name, DATE_TRUNC('month', check_timestamp)
                  HAVING SUM(
                    CASE
                      WHEN prev_stock IS NOT NULL AND stock_level < prev_stock
                      THEN (prev_stock - stock_level) * price
                      ELSE 0
                    END
                  ) > 0
                )
                SELECT
                  product_name,
                  month_start,
                  revenue
                FROM sales_calc
                ORDER BY month_start, product_name
            """)

            data = cur.fetchall()
            conn.close()

            # Группируем данные по месяцам
            monthly_data = {}
            for row in data:
                month_key = row['month_start'].strftime('%Y-%m')
                if month_key not in monthly_data:
                    monthly_data[month_key] = {}
                monthly_data[month_key][row['product_name']] = float(row['revenue'])

            # Получаем все уникальные товары
            all_products = set()
            for month_data in monthly_data.values():
                all_products.update(month_data.keys())

            # Создаем результат
            result = {
                'months': list(monthly_data.keys()),
                'products': list(all_products),
                'data': monthly_data
            }

            return result

    except Exception as e:
        current_app.logger.error(f"Ошибка получения данных выручки: {e}", exc_info=True)
        return {"error": "Ошибка получения данных выручки"}, 500

# --- API: экспорт графика продаж по месяцам в Excel ---
@bp.route('/api/export/chart/monthly/excel')
@login_required
def api_export_chart_monthly_excel():
    """API для экспорта графика продаж по месяцам в Excel"""
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, Reference
    import tempfile
    import os

    try:
        wb = Workbook()
        ws = wb.active
        ws.title = "Продажи по месяцам"

        conn = get_db_connection_webapp()
        if not conn:
            return {"error": "Ошибка подключения к БД"}, 500

        ws['A1'] = "График продаж по месяцам (последние 12 месяцев)"
        ws['A2'] = "Дата экспорта: " + datetime.now().strftime('%d.%m.%Y %H:%M')

        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                WITH months_series AS (
                  SELECT
                    DATE_TRUNC('month', CURRENT_DATE - INTERVAL '11 months' + (generate_series(0, 11) * INTERVAL '1 month')) AS month_start
                ),
                sh AS (
                  SELECT
                    p.id AS product_id,
                    p.name as product_name,
                    sh.check_timestamp,
                    sh.stock_level,
                    LAG(sh.stock_level) OVER (PARTITION BY sh.product_id ORDER BY sh.check_timestamp) AS prev_stock
                  FROM stock_history sh
                  JOIN products p ON p.id = sh.product_id
                  WHERE COALESCE(p.is_active, TRUE) = TRUE
                    AND sh.stock_level IS NOT NULL
                    AND sh.check_timestamp >= DATE_TRUNC('month', CURRENT_DATE - INTERVAL '11 months')
                ),
                sales_calc AS (
                  SELECT
                    product_id,
                    product_name,
                    DATE_TRUNC('month', check_timestamp) AS month_start,
                    SUM(
                      CASE
                        WHEN prev_stock IS NOT NULL AND stock_level < prev_stock
                        THEN prev_stock - stock_level
                        ELSE 0
                      END
                    ) AS sold_qty
                  FROM sh
                  GROUP BY product_id, product_name, DATE_TRUNC('month', check_timestamp)
                  HAVING SUM(
                    CASE
                      WHEN prev_stock IS NOT NULL AND stock_level < prev_stock
                      THEN prev_stock - stock_level
                      ELSE 0
                    END
                  ) > 0
                )
                SELECT
                  product_name,
                  month_start,
                  sold_qty
                FROM sales_calc
                ORDER BY product_name, month_start
            """)

            row = 4
            ws['A4'] = "Товар"
            ws['B4'] = "Месяц"
            ws['C4'] = "Продано (шт)"

            for row_data in cur.fetchall():
                row += 1
                ws[f'A{row}'] = row_data['product_name']
                ws[f'B{row}'] = row_data['month_start'].strftime('%Y-%m')
                ws[f'C{row}'] = int(row_data['sold_qty'])

        conn.close()

        chart = BarChart()
        chart.title = "Продажи по месяцам"
        chart.style = 10
        chart.y_axis.title = 'Количество проданных товаров'
        chart.x_axis.title = 'Месяц'

        data = Reference(ws, min_col=3, min_row=4, max_row=row)
        cats = Reference(ws, min_col=2, min_row=5, max_row=row)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)

        ws.add_chart(chart, "E2")

        with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp_file:
            wb.save(tmp_file.name)
            with open(tmp_file.name, 'rb') as f:
                file_data = f.read()
            os.unlink(tmp_file.name)

            from flask import Response
            return Response(
                file_data,
                mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                headers={
                    'Content-Disposition': f'attachment; filename=chipdip_monthly_{datetime.now().strftime("%Y%m%d_%H%M")}.xlsx'
                }
            )

    except Exception as e:
        current_app.logger.error(f"Ошибка экспорта графика месячных продаж в Excel: {e}", exc_info=True)
        return {"error": "Ошибка при создании Excel файла с графиком месячных продаж"}, 500


# --- API: экспорт графика остатков на складе в Excel ---
@bp.route('/api/export/chart/stock/excel')
@login_required
def api_export_chart_stock_excel():
    """API для экспорта графика остатков на складе в Excel для выбранного товара"""
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, LineChart, Reference
    import tempfile
    import os

    try:
        # Получаем ID товара из параметров
        product_id = request.args.get('product_id')
        if not product_id:
            return {"error": "Не указан ID товара"}, 400

        wb = Workbook()
        ws = wb.active
        ws.title = "Остатки на складе"

        conn = get_db_connection_webapp()
        if not conn:
            return {"error": "Ошибка подключения к БД"}, 500

        # Получаем название товара
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("SELECT name FROM products WHERE id = %s", (product_id,))
            product = cur.fetchone()
            if not product:
                return {"error": "Товар не найден"}, 404
            
            product_name = product['name']

        ws['A1'] = f"График остатков товара: {product_name}"
        ws['A2'] = "Дата экспорта: " + datetime.now().strftime('%d.%m.%Y %H:%M')

        # Получаем историю остатков для выбранного товара
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                SELECT 
                    check_timestamp,
                    stock_level
                FROM stock_history
                WHERE product_id = %s
                  AND stock_level IS NOT NULL
                ORDER BY check_timestamp ASC
                LIMIT 50
            """, (product_id,))

            row = 4
            ws['A4'] = "Дата проверки"
            ws['B4'] = "Остаток на складе"

            for row_data in cur.fetchall():
                row += 1
                ws[f'A{row}'] = row_data['check_timestamp'].strftime('%Y-%m-%d %H:%M')
                ws[f'B{row}'] = int(row_data['stock_level'])

        conn.close()

        chart = LineChart()
        chart.title = f"Остатки товара: {product_name}"
        chart.style = 10
        chart.y_axis.title = 'Остаток на складе'
        chart.x_axis.title = 'Дата проверки'

        data = Reference(ws, min_col=2, min_row=4, max_row=row)
        cats = Reference(ws, min_col=1, min_row=5, max_row=row)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)

        ws.add_chart(chart, "D2")

        with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp_file:
            wb.save(tmp_file.name)
            with open(tmp_file.name, 'rb') as f:
                file_data = f.read()
            os.unlink(tmp_file.name)

            from flask import Response
            return Response(
                file_data,
                mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                headers={
                    'Content-Disposition': f'attachment; filename=chipdip_stock_{product_id}_{datetime.now().strftime("%Y%m%d_%H%M")}.xlsx'
                }
            )

    except Exception as e:
        current_app.logger.error(f"Ошибка экспорта графика остатков в Excel: {e}", exc_info=True)
        return {"error": "Ошибка при создании Excel файла с графиком остатков"}, 500


# --- API: экспорт графика стоимости товара в Excel ---
@bp.route('/api/export/chart/price/excel')
@login_required
def api_export_chart_price_excel():
    """API для экспорта графика стоимости товара в Excel для выбранного товара"""
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, LineChart, Reference
    import tempfile
    import os

    try:
        # Получаем ID товара из параметров
        product_id = request.args.get('product_id')
        if not product_id:
            return {"error": "Не указан ID товара"}, 400

        wb = Workbook()
        ws = wb.active
        ws.title = "Стоимость товара"

        conn = get_db_connection_webapp()
        if not conn:
            return {"error": "Ошибка подключения к БД"}, 500

        # Получаем название товара
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("SELECT name FROM products WHERE id = %s", (product_id,))
            product = cur.fetchone()
            if not product:
                return {"error": "Товар не найден"}, 404
            
            product_name = product['name']

        ws['A1'] = f"График стоимости товара: {product_name}"
        ws['A2'] = "Дата экспорта: " + datetime.now().strftime('%d.%m.%Y %H:%M')

        # Получаем историю цен для выбранного товара
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                SELECT 
                    check_timestamp,
                    price
                FROM stock_history
                WHERE product_id = %s
                  AND price IS NOT NULL
                ORDER BY check_timestamp ASC
                LIMIT 50
            """, (product_id,))

            row = 4
            ws['A4'] = "Дата проверки"
            ws['B4'] = "Цена (руб)"

            for row_data in cur.fetchall():
                row += 1
                ws[f'A{row}'] = row_data['check_timestamp'].strftime('%Y-%m-%d %H:%M')
                ws[f'B{row}'] = float(row_data['price'])

        conn.close()

        chart = LineChart()
        chart.title = f"Стоимость товара: {product_name}"
        chart.style = 10
        chart.y_axis.title = 'Цена (руб)'
        chart.x_axis.title = 'Дата проверки'

        data = Reference(ws, min_col=2, min_row=4, max_row=row)
        cats = Reference(ws, min_col=1, min_row=5, max_row=row)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)

        ws.add_chart(chart, "D2")

        with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp_file:
            wb.save(tmp_file.name)
            with open(tmp_file.name, 'rb') as f:
                file_data = f.read()
            os.unlink(tmp_file.name)

            from flask import Response
            return Response(
                file_data,
                mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                headers={
                    'Content-Disposition': f'attachment; filename=chipdip_price_{product_id}_{datetime.now().strftime("%Y%m%d_%H%M")}.xlsx'
                }
            )

    except Exception as e:
        current_app.logger.error(f"Ошибка экспорта графика цен в Excel: {e}", exc_info=True)
        return {"error": "Ошибка при создании Excel файла с графиком цен"}, 500


# --- API: статистика выручки по месяцам для выбранного товара ---
@bp.route('/api/stats/revenue/product-monthly')
@login_required
def api_stats_revenue_product_monthly():
    """API для получения данных выручки по месяцам для выбранного товара"""
    product_id = request.args.get('product_id')
    if not product_id:
        return {"error": "Не указан ID товара"}, 400

    conn = get_db_connection_webapp()
    if not conn:
        return {"error": "Ошибка подключения к БД"}, 500

    try:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            # Проверяем существование товара
            cur.execute("""
                SELECT name FROM products 
                WHERE id = %s AND COALESCE(is_active, TRUE) = TRUE
            """, (product_id,))
            product = cur.fetchone()
            if not product:
                return {"error": "Товар не найден"}, 404

            cur.execute("""
                WITH sh AS (
                  SELECT
                    sh.check_timestamp,
                    sh.stock_level,
                    sh.price,
                    LAG(sh.stock_level) OVER (ORDER BY sh.check_timestamp) AS prev_stock
                  FROM stock_history sh
                  WHERE sh.product_id = %s
                    AND sh.stock_level IS NOT NULL
                    AND sh.price IS NOT NULL
                    AND sh.check_timestamp >= DATE_TRUNC('month', CURRENT_DATE - INTERVAL '11 months')
                ),
                sales_calc AS (
                  SELECT
                    DATE_TRUNC('month', check_timestamp) AS month_start,
                    SUM(
                      CASE
                        WHEN prev_stock IS NOT NULL AND stock_level < prev_stock
                        THEN (prev_stock - stock_level) * price
                        ELSE 0
                      END
                    ) AS revenue
                  FROM sh
                  GROUP BY DATE_TRUNC('month', check_timestamp)
                  HAVING SUM(
                    CASE
                      WHEN prev_stock IS NOT NULL AND stock_level < prev_stock
                      THEN (prev_stock - stock_level) * price
                      ELSE 0
                    END
                  ) > 0
                )
                SELECT
                  month_start,
                  revenue
                FROM sales_calc
                ORDER BY month_start ASC
            """, (product_id,))

            data = cur.fetchall()
            conn.close()

            # Формируем результат
            months = []
            revenues = []
            for row in data:
                months.append(row['month_start'].strftime('%Y-%m'))
                revenues.append(float(row['revenue']))

            return {
                'product_name': product['name'],
                'months': months,
                'revenues': revenues
            }

    except Exception as e:
        current_app.logger.error(f"Ошибка получения данных выручки для товара: {e}", exc_info=True)
        return {"error": "Ошибка получения данных выручки для товара"}, 500

# --- API: экспорт графика продаж по месяцам для выбранного товара в Excel ---
@bp.route('/api/export/chart/product-monthly/excel')
@login_required
def api_export_chart_product_monthly_excel():
    """API для экспорта графика продаж по месяцам для выбранного товара в Excel"""
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, LineChart, Reference
    import tempfile
    import os

    try:
        # Получаем ID товара из параметров
        product_id = request.args.get('product_id')
        if not product_id:
            return {"error": "Не указан ID товара"}, 400

        wb = Workbook()
        ws = wb.active
        ws.title = "Продажи по месяцам"

        conn = get_db_connection_webapp()
        if not conn:
            return {"error": "Ошибка подключения к БД"}, 500

        # Получаем название товара
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("SELECT name FROM products WHERE id = %s", (product_id,))
            product = cur.fetchone()
            if not product:
                return {"error": "Товар не найден"}, 404
            
            product_name = product['name']

        ws['A1'] = f"Продажи по месяцам: {product_name}"
        ws['A2'] = "Дата экспорта: " + datetime.now().strftime('%d.%m.%Y %H:%M')

        # Получаем данные продаж по месяцам для выбранного товара
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                WITH months_series AS (
                  SELECT 
                    DATE_TRUNC('month', CURRENT_DATE - INTERVAL '11 months' + (generate_series(0, 11) * INTERVAL '1 month')) AS month_start
                ),
                sh AS (
                  SELECT
                    sh.check_timestamp,
                    sh.stock_level,
                    LAG(sh.stock_level) OVER (ORDER BY sh.check_timestamp) AS prev_stock
                  FROM stock_history sh
                  WHERE sh.product_id = %s
                    AND sh.stock_level IS NOT NULL
                    AND sh.check_timestamp >= DATE_TRUNC('month', CURRENT_DATE - INTERVAL '11 months')
                ),
                sales_calc AS (
                  SELECT
                    DATE_TRUNC('month', check_timestamp) AS month_start,
                    SUM(
                      CASE 
                        WHEN prev_stock IS NOT NULL AND stock_level < prev_stock 
                        THEN prev_stock - stock_level
                        ELSE 0
                      END
                    ) AS sold_qty
                  FROM sh
                  GROUP BY DATE_TRUNC('month', check_timestamp)
                  HAVING SUM(
                    CASE 
                      WHEN prev_stock IS NOT NULL AND stock_level < prev_stock 
                      THEN prev_stock - stock_level
                      ELSE 0
                    END
                  ) > 0
                )
                SELECT 
                  ms.month_start,
                  COALESCE(sc.sold_qty, 0) AS sold_qty
                FROM months_series ms
                LEFT JOIN sales_calc sc ON ms.month_start = sc.month_start
                ORDER BY ms.month_start ASC
            """, (product_id,))

            row = 4
            ws['A4'] = "Месяц"
            ws['B4'] = "Продано (шт)"

            for row_data in cur.fetchall():
                row += 1
                ws[f'A{row}'] = row_data['month_start'].strftime('%Y-%m')
                ws[f'B{row}'] = int(row_data['sold_qty'])

        conn.close()

        chart = LineChart()
        chart.title = f"Продажи по месяцам: {product_name}"
        chart.style = 10
        chart.y_axis.title = 'Количество проданных товаров'
        chart.x_axis.title = 'Месяц'

        data = Reference(ws, min_col=2, min_row=4, max_row=row)
        cats = Reference(ws, min_col=1, min_row=5, max_row=row)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)

        ws.add_chart(chart, "D2")

        with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp_file:
            wb.save(tmp_file.name)
            with open(tmp_file.name, 'rb') as f:
                file_data = f.read()
            os.unlink(tmp_file.name)

            from flask import Response
            return Response(
                file_data,
                mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                headers={
                    'Content-Disposition': f'attachment; filename=chipdip_product_monthly_{product_id}_{datetime.now().strftime("%Y%m%d_%H%M")}.xlsx'
                }
            )

    except Exception as e:
        current_app.logger.error(f"Ошибка экспорта графика продаж по месяцам в Excel: {e}", exc_info=True)
        return {"error": "Ошибка при создании Excel файла с графиком продаж по месяцам"}, 500


# --- API: экспорт графика в PowerPoint ---
@bp.route('/api/export/chart/powerpoint')
@login_required
def api_export_chart_powerpoint():
    """API для экспорта графика в PowerPoint"""
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from io import BytesIO
    import tempfile
    import os

    try:
        # Создаем PowerPoint презентацию
        prs = Presentation()

        # Добавляем титульный слайд
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "График продаж ChipDip"
        subtitle.text = f"Отчет за {datetime.now().strftime('%d.%m.%Y %H:%M')}"

        # Получаем данные для экспорта
        conn = get_db_connection_webapp()
        if not conn:
            return {"error": "Ошибка подключения к БД"}, 500

        # Добавляем слайд с графиком
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        title.text = "График продаж по месяцам"

        # Получаем данные продаж
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                WITH months_series AS (
                  SELECT
                    DATE_TRUNC('month', CURRENT_DATE - INTERVAL '11 months' + (generate_series(0, 11) * INTERVAL '1 month')) AS month_start
                ),
                sh AS (
                  SELECT
                    p.id AS product_id,
                    p.name as product_name,
                    sh.check_timestamp,
                    sh.stock_level,
                    LAG(sh.stock_level) OVER (PARTITION BY sh.product_id ORDER BY sh.check_timestamp) AS prev_stock
                  FROM stock_history sh
                  JOIN products p ON p.id = sh.product_id
                  WHERE COALESCE(p.is_active, TRUE) = TRUE
                    AND sh.stock_level IS NOT NULL
                    AND sh.check_timestamp >= DATE_TRUNC('month', CURRENT_DATE - INTERVAL '11 months')
                ),
                sales_calc AS (
                  SELECT
                    product_id,
                    product_name,
                    DATE_TRUNC('month', check_timestamp) AS month_start,
                    SUM(
                      CASE
                        WHEN prev_stock IS NOT NULL AND stock_level < prev_stock
                        THEN prev_stock - stock_level
                        ELSE 0
                      END
                    ) AS sold_qty
                  FROM sh
                  GROUP BY product_id, product_name, DATE_TRUNC('month', check_timestamp)
                  HAVING SUM(
                    CASE
                      WHEN prev_stock IS NOT NULL AND stock_level < prev_stock
                      THEN prev_stock - stock_level
                      ELSE 0
                    END
                  ) > 0
                )
                SELECT
                  product_name,
                  month_start,
                  sold_qty
                FROM sales_calc
                ORDER BY product_name, month_start
            """)

            # Создаем данные для графика
            rows = cur.fetchall()
            if rows:
                # Группируем данные по месяцам
                months_data = {}
                for row_data in rows:
                    month_key = row_data['month_start'].strftime('%Y-%m')
                    if month_key not in months_data:
                        months_data[month_key] = {}
                    months_data[month_key][row_data['product_name']] = int(row_data['sold_qty'])

                # Создаем график
                chart_data = CategoryChartData()
                
                # Добавляем категории (месяцы)
                chart_data.categories = list(months_data.keys())
                
                # Добавляем серии данных (товары)
                all_products = set()
                for month_data in months_data.values():
                    all_products.update(month_data.keys())
                
                for product in sorted(all_products):
                    values = []
                    for month in chart_data.categories:
                        values.append(months_data[month].get(product, 0))
                    chart_data.add_series(product, values)

                # Добавляем график на слайд
                x, y, cx, cy = 100, 200, 600, 400
                chart = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
                ).chart

                chart.chart_title.text_frame.text = "Продажи по месяцам"

        conn.close()

        # Создаем временный файл
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
            prs.save(tmp_file.name)

            # Читаем файл и отправляем пользователю
            with open(tmp_file.name, 'rb') as f:
                file_data = f.read()

            # Удаляем временный файл
            os.unlink(tmp_file.name)

            from flask import Response
            return Response(
                file_data,
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                headers={
                    'Content-Disposition': f'attachment; filename=chipdip_chart_{datetime.now().strftime("%Y%m%d_%H%M")}.pptx'
                }
            )

    except Exception as e:
        current_app.logger.error(f"Ошибка экспорта графика в PowerPoint: {e}", exc_info=True)
        return {"error": "Ошибка при создании PowerPoint файла с графиком"}, 500
